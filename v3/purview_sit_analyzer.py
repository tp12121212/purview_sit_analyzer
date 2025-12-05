import argparse
import json
import logging
import os
import re
from bisect import bisect_right
from collections import Counter
from pathlib import Path
from typing import Dict, Iterable, List, Optional, Sequence, Tuple

import openpyxl
import pandas as pd
import pdfplumber
import pytesseract
import requests
from docx import Document
from pdf2image import convert_from_path
from PIL import Image
from pptx import Presentation

# ---- Configure Logging ----
logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
logger = logging.getLogger(__name__)

DEFAULT_STOPWORDS_URL = "https://raw.githubusercontent.com/tp12121212/purview_sit_analyzer/refs/heads/main/english.txt"
DEFAULT_LOCAL_STOPWORDS = Path(__file__).resolve().parent.parent / "english.txt"


def load_stopwords(local_path: Optional[Path] = None, remote_url: Optional[str] = None) -> set:
    """Load stopwords from a local file first, then optionally from a URL."""
    stopwords: set = set()
    local_candidates: Sequence[Path] = []
    if local_path:
        local_candidates = [local_path]
    else:
        local_candidates = [DEFAULT_LOCAL_STOPWORDS, Path("english.txt")]

    for candidate in local_candidates:
        try:
            if candidate.is_file():
                stopwords = set(candidate.read_text(encoding="utf-8").split())
                logger.info("Loaded %d stop words from %s", len(stopwords), candidate)
                break
        except Exception as exc:
            logger.warning("Failed to read stopwords from %s: %s", candidate, exc)

    if not stopwords and remote_url:
        try:
            response = requests.get(remote_url, timeout=10)
            response.raise_for_status()
            stopwords = set(response.text.split())
            logger.info("Loaded %d stop words from %s", len(stopwords), remote_url)
        except Exception as exc:
            logger.warning("Failed to fetch stopwords from %s: %s", remote_url, exc)

    return stopwords


def load_regex_patterns(patterns_path: Path) -> List[Dict]:
    with patterns_path.open("r", encoding="utf-8") as fh:
        return json.load(fh)


def precompile_patterns(patterns: Iterable[Dict]) -> List[Tuple[str, str, re.Pattern]]:
    compiled = []
    for pattern in patterns:
        name = pattern.get("name")
        regex = pattern.get("pattern")
        try:
            compiled.append((name, regex, re.compile(regex)))
        except re.error as exc:
            logger.warning("Skipping invalid regex %s: %s", name, exc)
    return compiled


def extract_text_from_file(filepath: Path, ocr: bool = True) -> str:
    """Extract text from supported file types; optionally OCR images/PDFs."""
    ext = filepath.suffix.lower()
    text = ""
    try:
        if ext in [".txt", ".csv"]:
            text = filepath.read_text(encoding="utf-8", errors="ignore")
        elif ext == ".docx":
            doc = Document(filepath)
            text = "\n".join([p.text for p in doc.paragraphs])
        elif ext == ".xlsx":
            wb = openpyxl.load_workbook(filepath)
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws.iter_rows(values_only=True):
                    text += " ".join([str(cell) for cell in row if cell]) + "\n"
        elif ext == ".pptx":
            prs = Presentation(filepath)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text += shape.text + "\n"
        elif ext == ".pdf":
            with pdfplumber.open(filepath) as pdf:
                for page in pdf.pages:
                    text += page.extract_text() or ""
            if not text.strip() and ocr:
                images = convert_from_path(filepath)
                for img in images:
                    text += pytesseract.image_to_string(img)
        elif ext in [".png", ".jpeg", ".jpg", ".gif"]:
            if ocr:
                img = Image.open(filepath)
                text = pytesseract.image_to_string(img)
    except Exception as exc:
        logger.error("Error extracting text from %s: %s", filepath, exc)
    return text


def build_line_index(text: str) -> List[int]:
    """Return list of character offsets for the start of each line."""
    offsets = [0]
    running = 0
    for line in text.splitlines(keepends=True):
        running += len(line)
        offsets.append(running)
    return offsets


def line_number_from_offset(line_offsets: Sequence[int], offset: int) -> int:
    """Map a character offset to a 1-based line number."""
    return bisect_right(line_offsets, offset)


def find_regex_matches(
    text: str, patterns: Sequence[Tuple[str, str, re.Pattern]]
) -> List[Dict]:
    """Collect regex detections with first-line metadata."""
    results = []
    line_offsets = build_line_index(text)
    for name, raw_pattern, compiled in patterns:
        matches = list(compiled.finditer(text))
        if not matches:
            continue
        first_line = line_number_from_offset(line_offsets, matches[0].start())
        results.append(
            {
                "Regex_Name": name,
                "Regex_Pattern": raw_pattern,
                "Detected_Text": matches[0].group(0),
                "Line_Number": first_line,
                "Match_Count": len(matches),
            }
        )
    return results


def normalize_for_stopwords(token: str) -> str:
    """Trim outer punctuation for stopword/length checks without touching interiors (e.g., emails)."""
    return token.strip(".,;:!?\"'()[]{}").lower()


def contiguous_phrases(
    text: str,
    stop_words: set,
    min_keyword_length: int,
    min_phrase_words: int,
    max_phrase_words: int,
) -> Counter:
    """
    Build phrases only from tokens that are next to each other in the original text
    and separated by literal spaces (no skipping stopwords or merging across lines).
    """
    tokens_with_separators: List[Tuple[str, str]] = []
    prev_end = 0
    for match in re.finditer(r"\S+", text):
        separator = text[prev_end : match.start()]
        tokens_with_separators.append((match.group(), separator))
        prev_end = match.end()

    phrases = Counter()
    run: List[str] = []
    for token, separator in tokens_with_separators:
        normalized = normalize_for_stopwords(token)
        token_ok = normalized and len(normalized) >= min_keyword_length and normalized not in stop_words
        separator_ok = separator == "" or set(separator) <= {" "}

        if token_ok and (not run or separator_ok):
            run.append(token)
        else:
            run.clear()
            if token_ok:
                run.append(token)

        if run:
            for length in range(min_phrase_words, max_phrase_words + 1):
                if len(run) >= length:
                    phrase = " ".join(run[-length:])
                    phrases[phrase] += 1
    return phrases


def analyze_files(
    input_folder: Path,
    compiled_patterns: Sequence[Tuple[str, str, re.Pattern]],
    stop_words: set,
    min_keyword_length: int,
    min_phrase_words: int,
    max_phrase_words: int,
    ocr: bool,
) -> Tuple[List[Dict], List[Dict]]:
    regex_results: List[Dict] = []
    keyword_results: List[Dict] = []

    for root, _, files in os.walk(input_folder):
        for file in files:
            filepath = Path(root) / file
            logger.info("Processing file: %s", filepath)
            text = extract_text_from_file(filepath, ocr=ocr)

            regex_matches = find_regex_matches(text, compiled_patterns)
            for match in regex_matches:
                match["File_Name"] = file
                regex_results.append(match)

            phrase_counts = contiguous_phrases(
                text,
                stop_words=stop_words,
                min_keyword_length=min_keyword_length,
                min_phrase_words=min_phrase_words,
                max_phrase_words=max_phrase_words,
            )
            for keyword, count in phrase_counts.items():
                keyword_results.append({"File_Name": file, "Keyword": keyword, "Count": count})

    return regex_results, keyword_results


def main() -> None:
    parser = argparse.ArgumentParser(description="Purview SIT Analyzer")
    parser.add_argument("--input_folder", required=True, help="Folder containing files to analyze")
    parser.add_argument("--output_folder", required=True, help="Folder to save output CSV files")
    parser.add_argument("--min_keyword_length", type=int, default=4, help="Minimum length of keyword")
    parser.add_argument("--min_phrase_words", type=int, default=2, help="Minimum words in phrase")
    parser.add_argument("--max_phrase_words", type=int, default=4, help="Maximum words in phrase")
    parser.add_argument("--stopwords_url", default=None, help="Optional stopwords URL override")
    parser.add_argument("--stopwords_path", type=Path, default=None, help="Optional local stopwords path")
    parser.add_argument("--no_ocr", action="store_true", help="Skip OCR for PDFs/images")
    args = parser.parse_args()

    input_folder = Path(args.input_folder)
    output_folder = Path(args.output_folder)
    output_folder.mkdir(parents=True, exist_ok=True)

    stop_words = load_stopwords(local_path=args.stopwords_path, remote_url=args.stopwords_url or DEFAULT_STOPWORDS_URL)
    regex_patterns = load_regex_patterns(Path(__file__).resolve().parent / "regex_patterns.json")
    compiled_patterns = precompile_patterns(regex_patterns)

    regex_results, keyword_results = analyze_files(
        input_folder=input_folder,
        compiled_patterns=compiled_patterns,
        stop_words=stop_words,
        min_keyword_length=args.min_keyword_length,
        min_phrase_words=args.min_phrase_words,
        max_phrase_words=args.max_phrase_words,
        ocr=not args.no_ocr,
    )

    regex_df = pd.DataFrame(regex_results)
    regex_df.to_csv(output_folder / "regex_detections.csv", index=False)

    keywords_df = pd.DataFrame(keyword_results)
    keywords_df.to_csv(output_folder / "keywords.csv", index=False)


if __name__ == "__main__":
    main()
