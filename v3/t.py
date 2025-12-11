import argparse
import csv
import json
import logging
import os
import re
from bisect import bisect_right
from pathlib import Path
from typing import Dict, Iterable, List, Optional, Sequence, Tuple
from collections import defaultdict

import pandas as pd
import requests
import pdfplumber
import pytesseract
import openpyxl
from pdf2image import convert_from_path
from PIL import Image
from docx import Document
from pptx import Presentation

logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
logger = logging.getLogger(__name__)

DEFAULT_REGEX_PATH = Path(__file__).resolve().parent / "regex_patterns.json"
DEFAULT_LEXICON_PATH = Path(__file__).resolve().parent.parent / "wordlists" / "lexicon_dict" / "lexicon_latest.csv"
DEFAULT_STOPWORDS_URL = "https://raw.githubusercontent.com/tp12121212/purview_sit_analyzer/refs/heads/main/english.txt"
DEFAULT_LOCAL_STOPWORDS = Path(__file__).resolve().parent.parent / "english.txt"


def load_stopwords(local_path: Optional[Path] = None, remote_url: Optional[str] = None) -> set:
    stopwords: set = set()
    candidates = [local_path] if local_path else [DEFAULT_LOCAL_STOPWORDS, Path("english.txt")]
    for cand in candidates:
        try:
            if cand and cand.is_file():
                stopwords = set(cand.read_text(encoding="utf-8").split())
                logger.info("Loaded %d stop words from %s", len(stopwords), cand)
                break
        except Exception as exc:
            logger.warning("Failed to read stopwords from %s: %s", cand, exc)
    if not stopwords and remote_url:
        try:
            resp = requests.get(remote_url, timeout=10)
            resp.raise_for_status()
            stopwords = set(resp.text.split())
            logger.info("Loaded %d stop words from %s", len(stopwords), remote_url)
        except Exception as exc:
            logger.warning("Failed to fetch stopwords from %s: %s", remote_url, exc)
    return stopwords


def load_regex_patterns(patterns_path: Path) -> List[Dict]:
    with patterns_path.open("r", encoding="utf-8") as fh:
        return json.load(fh)


def precompile_patterns(patterns: Iterable[Dict]) -> List[Tuple[str, str, re.Pattern]]:
    compiled = []
    for pattern in patterns:
        name = pattern.get("name")
        regex = pattern.get("pattern")
        try:
            compiled.append((name, regex, re.compile(regex)))
        except re.error as exc:
            logger.warning("Skipping invalid regex %s: %s", name, exc)
    return compiled


def load_lexicon(path: Path) -> Dict[str, float]:
    """
    Load lexicon_latest.csv with columns: keyword,length.
    Returns mapping keyword(lower) -> length (float).
    """
    lex: Dict[str, float] = {}
    if not path.is_file():
        logger.warning("Lexicon file not found at %s; no keyword pairs will be detected.", path)
        return lex
    try:
        df = pd.read_csv(path)
        for _, row in df.iterrows():
            kw = str(row.get("keyword", "")).strip().lower()
            if not kw:
                continue
            try:
                length_val = float(row.get("length", 0))
            except Exception:
                length_val = 0.0
            lex[kw] = length_val
        logger.info("Loaded %d lexicon entries from %s", len(lex), path)
    except Exception as exc:
        logger.error("Failed to load lexicon from %s: %s", path, exc)
    return lex


def build_line_index(text: str) -> List[int]:
    offsets = [0]
    running = 0
    for line in text.splitlines(keepends=True):
        running += len(line)
        offsets.append(running)
    return offsets


def line_number_from_offset(line_offsets: Sequence[int], offset: int) -> int:
    return bisect_right(line_offsets, offset)


def find_regex_matches(text: str, patterns: Sequence[Tuple[str, str, re.Pattern]]) -> List[Dict]:
    results = []
    line_offsets = build_line_index(text)
    for name, raw_pattern, compiled in patterns:
        matches = list(compiled.finditer(text))
        if not matches:
            continue
        first_line = line_number_from_offset(line_offsets, matches[0].start())
        results.append(
            {
                "Regex_Name": name,
                "Regex_Pattern": raw_pattern,
                "Detected_Text": matches[0].group(0),
                "Line_Number": first_line,
                "Match_Count": len(matches),
            }
        )
    return results


def detect_adjacent_pairs(text: str, lexicon: Dict[str, float], stop_words: set) -> Dict[str, Tuple[int, float]]:
    """
    Find adjacent two-token pairs where both tokens are in the lexicon (case-insensitive)
    and not in the stopword list. Tokens are split on whitespace only—no punctuation stripping.
    Returns mapping pair -> (count, total_weight).
    """
    results: Dict[str, Tuple[int, float]] = defaultdict(lambda: (0, 0.0))
    if not lexicon:
        return results

    tokens = text.split()
    for i in range(len(tokens) - 1):
        raw1, raw2 = tokens[i], tokens[i + 1]
        k1, k2 = raw1.lower(), raw2.lower()
        if k1 in stop_words or k2 in stop_words:
            continue
        if k1 in lexicon and k2 in lexicon:
            pair = f"{raw1} {raw2}"
            w = lexicon[k1] + lexicon[k2]
            c, wt = results[pair]
            results[pair] = (c + 1, wt + w)
    return results


def read_text_file(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8", errors="ignore")
    except Exception as exc:
        logger.error("Failed to read %s: %s", path, exc)
        return ""


def extract_text_from_file(filepath: Path, ocr: bool = True) -> str:
    ext = filepath.suffix.lower()
    text = ""
    try:
        if ext in [".txt", ".csv"]:
            text = filepath.read_text(encoding="utf-8", errors="ignore")
        elif ext == ".docx":
            doc = Document(filepath)
            text = "\n".join([p.text for p in doc.paragraphs])
        elif ext == ".xlsx":
            wb = openpyxl.load_workbook(filepath)
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws.iter_rows(values_only=True):
                    text += " ".join([str(cell) for cell in row if cell]) + "\n"
        elif ext == ".pptx":
            prs = Presentation(filepath)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text += shape.text + "\n"
        elif ext == ".pdf":
            with pdfplumber.open(filepath) as pdf:
                for page in pdf.pages:
                    text += page.extract_text() or ""
            if not text.strip() and ocr:
                images = convert_from_path(filepath)
                for img in images:
                    text += pytesseract.image_to_string(img)
        elif ext in [".png", ".jpeg", ".jpg", ".gif", ".tiff", ".tif", ".bmp"]:
            if ocr:
                img = Image.open(filepath)
                text = pytesseract.image_to_string(img)
    except Exception as exc:
        logger.error("Error extracting text from %s: %s", filepath, exc)
    return text


def main() -> None:
    parser = argparse.ArgumentParser(description="Regex + adjacent lexicon pair detector")
    parser.add_argument("--input_folder", required=True, help="Folder containing files to analyze")
    parser.add_argument("--output_folder", required=True, help="Folder to save output CSV files")
    parser.add_argument("--regex_path", type=Path, default=None, help="Path to regex_patterns.json")
    parser.add_argument("--lexicon_path", type=Path, default=None, help="Path to lexicon_latest.csv")
    parser.add_argument("--stopwords_path", type=Path, default=None, help="Optional local stopwords path")
    parser.add_argument("--stopwords_url", default=None, help="Optional stopwords URL override")
    parser.add_argument("--no_ocr", action="store_true", help="Skip OCR for PDFs/images")
    args = parser.parse_args()

    input_folder = Path(args.input_folder)
    output_folder = Path(args.output_folder)
    output_folder.mkdir(parents=True, exist_ok=True)

    stop_words = load_stopwords(local_path=args.stopwords_path, remote_url=args.stopwords_url or DEFAULT_STOPWORDS_URL)
    lexicon = load_lexicon(args.lexicon_path or DEFAULT_LEXICON_PATH)
    regex_patterns = load_regex_patterns(args.regex_path or DEFAULT_REGEX_PATH)
    compiled_patterns = precompile_patterns(regex_patterns)

    regex_rows = []
    keyword_rows = []

    for root, _, files in os.walk(input_folder):
        for fname in files:
            fpath = Path(root) / fname
            text = extract_text_from_file(fpath, ocr=not args.no_ocr)

            # Debug output: raw and normalized snippets
            raw_snippet = text[:2000]
            normalized_tokens = " ".join(text.lower().split())[:2000]
            print(f"\n=== RAW TEXT: {fpath} ===\n{raw_snippet}\n")
            print(f"=== NORMALIZED TOKENS: {fpath} ===\n{normalized_tokens}\n")

            # Regex detection
            regex_matches = find_regex_matches(text, compiled_patterns)
            for match in regex_matches:
                match["File_Name"] = fname
                match["File_Path"] = str(fpath)
                regex_rows.append(match)

            # Adjacent lexicon pairs
            pairs = detect_adjacent_pairs(text, lexicon, stop_words)
            for pair, (count, weight) in pairs.items():
                keyword_rows.append(
                    {
                        "File_Name": fname,
                        "File_Path": str(fpath),
                        "Keyword": pair,
                        "Count": count,
                        "Weight": weight,
                    }
                )

    pd.DataFrame(regex_rows).to_csv(output_folder / "regex_detections.csv", index=False)

    if keyword_rows:
        keyword_rows = sorted(keyword_rows, key=lambda r: (r["Weight"], r["Count"]), reverse=True)
    pd.DataFrame(keyword_rows).to_csv(output_folder / "keywords.csv", index=False)

    logger.info("Wrote %d regex rows and %d keyword rows to %s", len(regex_rows), len(keyword_rows), output_folder)


if __name__ == "__main__":
    main()
